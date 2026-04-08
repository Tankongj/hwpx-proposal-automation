#!/usr/bin/env python3
"""
extract_ppt.py — Extract ordered structure and text from PPT/PPTX files.

Reads a PowerPoint file and extracts text, table data, and structure by slide.
Outputs to a JSON file to serve as a structured input for LLM semantic mapping.

Usage:
    python extract_ppt.py input.pptx [output.json]
"""

import sys
import os
import json
import argparse
from pptx import Presentation

def parse_ppt(filepath: str, output_path: str = None) -> dict:
    """Extract slide text, tables, and shape structures into a dict."""
    if not os.path.exists(filepath):
        print(f"Error: File not found: {filepath}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(filepath)
    result = {"filename": os.path.basename(filepath), "slides": []}

    for i, slide in enumerate(prs.slides):
        slide_layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        slide_data = {
            "slide_index": i + 1,
            "layout": slide_layout_name,
            "content": []
        }
        
        for shape in slide.shapes:
            shape_info = {"type": str(shape.shape_type)}
            if hasattr(shape, "name"):
                shape_info["name"] = shape.name

            # Extract Text
            if shape.has_text_frame:
                paragraphs = []
                for p in shape.text_frame.paragraphs:
                    text = p.text.strip()
                    if text:
                        paragraphs.append({"level": p.level, "text": text})
                if paragraphs:
                    shape_info["paragraphs"] = paragraphs
                    shape_info["full_text"] = shape.text

            # Extract Table
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                shape_info["table"] = table_data
                
            # Log placeholder presence
            if shape.is_placeholder:
                shape_info["is_placeholder"] = True

            # Extract Image (we don't save blobs to avoid massive JSON, just flag it)
            if hasattr(shape, "image"):
                shape_info["has_image"] = True
                shape_info["image_type"] = shape.image.content_type
            
            # Add shape to content if it actually contains text, table or image
            if "paragraphs" in shape_info or "table" in shape_info or "has_image" in shape_info:
                slide_data["content"].append(shape_info)
                
        result["slides"].append(slide_data)

    # Output handling
    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        print(f"✅ Saved PPT structure to: {output_path}")
    else:
        print(json.dumps(result, ensure_ascii=False, indent=2))

    return result

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Extract structured text and tables from PPT/PPTX')
    parser.add_argument('pptx_path', help='Path to PPTX file')
    parser.add_argument('output', nargs='?', default=None, help='Output JSON file path (optional)')

    args = parser.parse_args()
    parse_ppt(args.pptx_path, args.output)
